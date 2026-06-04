#!/usr/bin/env python3
"""Generate compact Chinese review notes from EC508 NLP lecture files."""

from __future__ import annotations

import re
import sys
from pathlib import Path
from typing import Iterable

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "EC508_NLP_lecture_notes_zh.docx"
SUPPORTED = {".pdf", ".pptx", ".docx", ".txt", ".md"}
ACCENT = RGBColor(31, 78, 121)
FONT_FAMILY = "Calibri"
BODY_SIZE = 10.5


LECTURE_NOTES = {
    "Lecture 02 -- What is NLP.pdf": {
        "title": "Lecture 02: 什么是自然语言处理 / What is NLP",
        "overview": (
            "本讲建立自然语言处理（Natural Language Processing, NLP）的整体地图。"
            "NLP 是人工智能（AI）的核心组成部分，也与机器学习（Machine Learning）、"
            "深度学习（Deep Learning）、语言学和信息检索密切相关。课程以 Turing Test、"
            "聊天机器人和 NLP 历史为背景，重点梳理文本处理、分类、信息检索、知识图谱、"
            "文本生成、问答、推理以及多模态转换等任务。"
        ),
        "points": [
            "自然语言处理（NLP）关注计算机对人类语言的分析、理解、生成与转换；语言是人类活动的基础，因此 NLP 的应用面非常广。",
            "图灵测试（Turing Test）用文本对话检验机器能否表现出近似人类的智能；ELIZA、Loebner Prize 与 ChatGPT 体现了聊天机器人的演进。",
            "NLP 的方法经历了形式语言与逻辑方法、概率方法、有限状态模型，以及融合机器学习和深度学习的阶段。",
            "文本预处理包括拼写纠正、规范化（normalization）、分词（tokenization）、词干提取（stemming）、词形还原（lemmatization）、词性标注（POS tagging）和词义消歧（word sense disambiguation）。",
            "文本分类（text classification）可使用规则、传统机器学习、深度学习（如 BERT）或混合方法；典型应用包括垃圾邮件、主题、紧急程度、意图和情感分析。",
            "实体与关系抽取包括命名实体识别（NER）、指代消解（coreference resolution）和关系抽取（relation extraction），结果可组织为知识图谱（knowledge graph）。",
            "生成类任务包括机器翻译（machine translation）、文本生成（text generation）、摘要（summarization）、问答（question answering）、文本到图像（text-to-image）和语音文本互转。",
        ],
        "terms": [
            ("Natural Language Processing (NLP)", "自然语言处理：让计算机处理、理解和生成人类语言。"),
            ("Turing Test", "图灵测试：通过文本问答判断机器是否能模拟人类智能。"),
            ("text classification", "文本分类：为句子或文档分配类别标签。"),
            ("sentiment analysis", "情感分析：判断文本的态度、情绪或倾向程度。"),
            ("named entity recognition (NER)", "命名实体识别：识别人名、地点、组织等实体及其类型。"),
            ("knowledge graph", "知识图谱：用实体节点与语义关系边组织信息。"),
            ("question answering (QA)", "问答：根据文本查询返回答案，可分为抽取式和生成式。"),
        ],
        "methods": [
            "分类系统常见路线：规则方法 -> 传统 Machine Learning（Naive Bayes、Logistic Regression、SVM）-> Deep Learning（如 BERT）-> 混合方法。",
            "实体信息管线：NER 识别实体 -> coreference resolution 合并同一实体的不同提法 -> relation extraction 提取关系 -> 构建 knowledge graph。",
            "问答系统区分抽取式问答（extractive QA）与生成式问答（generative QA）：前者从文档定位片段，后者生成自然语言回答。",
        ],
        "review": [
            "能按“预处理、分类、检索、抽取、生成、问答、推理、多模态”归类 NLP 任务。",
            "理解 Turing Test 的设置、意义与局限，并能说明聊天机器人为何是 NLP 的综合应用。",
            "区分 NER、coreference resolution、relation extraction 与 knowledge graph 的作用。",
        ],
    },
    "Lecture 03 -- Basic Notions and Text Normalization.pdf": {
        "title": "Lecture 03: 基础概念与文本规范化 / Basic Notions and Text Normalization",
        "overview": (
            "本讲讨论原始文本如何转化为可供模型使用的数据。文本最底层是有限字符集合上的序列，"
            "但实际 NLP 通常需要将其组织成 token、句子、文档和 corpus。重点是文本数据准备流程、"
            "分词（tokenization）、规范化（normalization）、词干提取（stemming）、词形还原"
            "（lemmatization）和句子切分（sentence segmentation）。"
        ),
        "points": [
            "书写系统可粗略分为 alphabet、syllabary 和 logography；NLP 的统一视角是把文本看作有限符号集合上的序列，常由 Unicode 表示。",
            "层级结构依次为 character、word、token、sentence、paragraph、document、corpus。token 是经过预处理后供算法使用的单位。",
            "真实项目通常要从抓取、OCR、语音转文本、格式转换开始，再做清洗、tokenization、词形规范化和句子切分。",
            "简单按空格或标点切词并不可靠：缩写、金额、日期、URL、hashtag、email、clitic 和 multiword expression (MWE) 都需要特殊处理。",
            "不同语言的 tokenization 难度不同；中文等语言不总用空格标记词边界，可用规则、Hidden Markov Model 或神经网络方法。",
            "规范化（normalization）把大小写、拼写、缩写、连字符等表面差异统一为标准形式。",
            "词干提取（stemming）用规则截取词干；词形还原（lemmatization）返回词典中的 lemma，通常需要更强的语言知识。",
            "句子切分要判断句点是句界还是缩写、数字的一部分；通常先 tokenization，再结合规则或机器学习分类。",
        ],
        "terms": [
            ("corpus", "语料库：用于训练、评估或分析的一组文档。"),
            ("tokenization", "分词：将字符序列切分为可处理的 token。"),
            ("normalization", "规范化：将文本表面形式转换成一致标准。"),
            ("stemming", "词干提取：按规则去除词缀，得到近似词干。"),
            ("lemmatization", "词形还原：将变形词映射到词典中的 lemma。"),
            ("multiword expression (MWE)", "多词表达：语义上应作为整体处理的词组，如 New York。"),
            ("sentence segmentation", "句子切分：识别文本中的句子边界。"),
        ],
        "methods": [
            "典型数据准备流程：raw data -> cleaning / de-noising -> character sequence -> tokenization -> normalization / stemming / lemmatization -> sentence segmentation -> corpus。",
            "规则式 tokenization 可使用 regular expression，但必须为 U.S.A.、$12.40、82%、URL、email 等模式设置例外。",
            "Porter Stemmer 是多阶段规则系统；lemmatization 则尝试将 am / are / is 统一为 be，将 reading 统一为 read。",
        ],
        "review": [
            "区分 word 与 token，区分 stemming 与 lemmatization。",
            "解释为什么“按空格切分”与“删除所有标点”会失败，并能举出至少三类反例。",
            "记住文本准备管线以及 tokenization 和 sentence segmentation 的先后关系。",
        ],
    },
    "Lecture 05 -- Language Models, Smoothing, Perplexity.pdf": {
        "title": "Lecture 05: 语言模型、平滑与困惑度 / Language Models, Smoothing, Perplexity",
        "overview": (
            "本讲复习生成式 n-gram language model，并讨论语言模型的评估。语言模型为 token 序列"
            "分配概率；训练语料只是无限语言空间中的有限样本。课程区分外部评估与内部评估，重点推导"
            "困惑度（perplexity），并指出零概率会使 perplexity 无法计算，从而引出 smoothing。"
        ),
        "points": [
            "n-gram language model 根据前 N-1 个 token 估计下一个 token 的条件概率，常用 N 为 2 到 4。",
            "生成文本时，先加入句界标记 <s> 与 </s>，再从匹配当前上下文的 n-gram 分布中逐步采样。",
            "语言模型（language model）本质上是 token 序列上的概率分布；有限训练集要近似无限大的语言空间。",
            "样本质量取决于规模与代表性：新闻语料未必适合训练聊天机器人，通用模型需要来源多样的 corpus。",
            "外部评估（extrinsic evaluation）在真实任务中比较准确率；内部评估（intrinsic evaluation）通常使用 train/test split。",
            "困惑度（perplexity）衡量模型对测试文本的不确定性：低 perplexity 较好，高 perplexity 较差。",
            "长句的联合概率天然更小，因此要用几何平均进行长度归一化；计算中常转到 log space 以避免 underflow。",
            "未见过的 n-gram 会得到 0 概率，使整句概率为 0；需要 smoothing 为未见事件分配非零概率。",
        ],
        "terms": [
            ("language model (LM)", "语言模型：对 token 序列分配概率的模型。"),
            ("n-gram", "n 元语法：由连续 N 个 token 构成的局部上下文。"),
            ("train/test split", "训练测试划分：在训练集建模，在测试集评估泛化能力。"),
            ("perplexity", "困惑度：衡量语言模型对测试文本不确定性的指标，越低越好。"),
            ("log space", "对数空间：将概率乘法变成加法，减少数值下溢。"),
            ("smoothing", "平滑：给未见事件分配概率质量，避免零概率。"),
        ],
        "methods": [
            "链式概率：P(W) = P(w1) * P(w2 | w1) * ... * P(wn | w1...w(n-1))。",
            "n-gram 近似：P(wn | w1...w(n-1)) ≈ P(wn | w(n-N+1)...w(n-1))。",
            "长度归一化概率：Pnorm(W) = P(W)^(1/n)。困惑度：PP(W) = 1 / Pnorm(W) = P(W)^(-1/n)。",
            "对数形式：log P(W) = sum_i log P(wi | context_i)；PP(W) = exp(-(1/n) * log P(W))。",
        ],
        "review": [
            "会解释 n-gram 的生成过程、Markov 假设和句界标记的作用。",
            "会推导 perplexity，说明为什么要做长度归一化和 log-space 计算。",
            "理解零概率问题：任何一个条件概率为 0，整句概率即为 0，perplexity 失效。",
        ],
    },
    "Lecture 06 -- Introduction to Term Vector Models.pdf": {
        "title": "Lecture 06: 平滑与词项向量模型 / Smoothing and Term Vector Models",
        "overview": (
            "本讲先完成 n-gram smoothing，再引入词项向量模型。平滑方法包括 Add-one estimation、"
            "backoff、interpolation 和 Kneser-Ney。随后用词袋模型（bag-of-words / BOW）表示文档，"
            "通过 TF-IDF 强调更有区分度的词，并用余弦相似度（cosine similarity）衡量向量接近程度。"
        ),
        "points": [
            "稀疏训练数据会造成 zero-probability n-gram；平滑（smoothing）的直觉是从已见事件中转移一部分概率质量给未见事件。",
            "Add-one estimation 又称 Laplace smoothing：每个词的计数都加 1。它简单，但通常不适合高质量 language model。",
            "回退（backoff）在高阶 n-gram 证据不足时减少左侧上下文；插值（interpolation）则对不同阶概率做加权求和。",
            "Stupid Backoff 递归使用固定折扣因子，如 0.4；更高级的方法包括 interpolated Kneser-Ney。",
            "词袋模型（bag-of-words / BOW）忽略词序，用 vocabulary 上的词频向量表示文本。",
            "词频（term frequency, TF）表示 term t 在 document d 中出现的次数；文档频率（document frequency, DF）表示包含 t 的文档数。",
            "逆文档频率（inverse document frequency, IDF）让稀有词获得较大权重；TF-IDF 同时考虑文档内频率与语料库中的稀有程度。",
            "欧氏距离受文档长度影响；余弦相似度（cosine similarity）比较向量夹角，更适合比较文本方向。",
        ],
        "terms": [
            ("Add-one / Laplace smoothing", "加一平滑：给每个候选词的计数加 1。"),
            ("backoff", "回退：高阶 n-gram 证据不足时使用更短上下文。"),
            ("interpolation", "插值：加权组合 unigram、bigram、trigram 等概率。"),
            ("bag-of-words (BOW)", "词袋模型：忽略词序，以词频向量表示文本。"),
            ("TF-IDF", "词频-逆文档频率：提高文档内常见但跨文档稀有词的权重。"),
            ("cosine similarity", "余弦相似度：用向量夹角的余弦度量文本相似性。"),
        ],
        "methods": [
            "最大似然 bigram：PMLE(wi | w(i-1)) = c(w(i-1), wi) / c(w(i-1))。",
            "加一平滑：PAdd-1(wi | w(i-1)) = (c(w(i-1), wi) + 1) / (c(w(i-1)) + V)，其中 V 为 vocabulary 大小。",
            "逆文档频率：IDF(t) = log(N / DFt)，其中 N 为文档总数，DFt 为包含 term t 的文档数。",
            "TF-IDF：w(t,d) = TF(t,d) * IDF(t)。课件还给出常见变体：w(t,d) = (1 + log10 TF(t,d)) * log10(N / DFt)。",
            "余弦相似度：cos(q,d) = (q · d) / (||q|| * ||d||)。值越接近 1，方向越相似。",
        ],
        "review": [
            "会比较 Add-one、backoff、interpolation 与 Kneser-Ney 的目的和基本思路。",
            "会计算 TF、DF、IDF、TF-IDF，并解释稀有词为何更有信息量。",
            "能说明 cosine similarity 为什么通常比 Euclidean distance 更适合文档比较。",
        ],
    },
    "Lecture 07 -- Vector Models, Word and Text Embeddings.pdf": {
        "title": "Lecture 07: 向量模型、PCA 与词嵌入 / Vector Models and Embeddings",
        "overview": (
            "本讲从高维稀疏 TF / TF-IDF 向量过渡到短而稠密的嵌入（embedding）。"
            "课程先用余弦相似度和 PCA 分析文本向量，再介绍分布式语义（distributional semantics）、"
            "word2vec、类比推理以及文档嵌入。"
        ),
        "points": [
            "TF BOW 与 TF-IDF 位于 |V| 维空间：term 是坐标轴，文本是向量。该空间通常高维且稀疏。",
            "余弦相似度（cosine similarity）忽略向量长度，关注方向；其理论范围为 [-1, 1]。",
            "主成分分析（PCA）通过奇异值分解（SVD）找到数据方差最大的方向，可把高维文本向量投影到二维用于观察。",
            "显式词频向量效率低、泛化弱，难以表达同义词与一词多义；短而稠密的 vector 更适合作为机器学习特征。",
            "分布式语义（distributional semantics）的核心是“词的意义来自其上下文”：相似语境中的词往往语义相近。",
            "word embedding 是短而稠密的词或文本向量；word2vec 包括 skip-gram 和 CBOW 等训练思路。",
            "skip-gram 上下文可理解为中心词前后窗口中的邻近词；训练时用共现 pair 作为正例、非共现 pair 作为负例调整 embedding。",
            "embedding 支持向量算术和简单类比推理，例如 tree - apple + grape，并通过 cosine similarity 搜索最近词。",
            "静态 embedding 为词分配固定向量；ELMo、BERT 等 contextual embedding 会根据上下文为同一个词生成不同向量。",
            "简单 document embedding 可以取所有 word embedding 的和，但更复杂任务通常会用更强的编码方法。",
        ],
        "terms": [
            ("PCA", "主成分分析：保留最大方差方向的降维方法。"),
            ("SVD", "奇异值分解：可用于分解 term-document matrix，并获得主成分。"),
            ("distributional semantics", "分布式语义：通过词的上下文使用方式刻画意义。"),
            ("word embedding", "词嵌入：词在低维稠密向量空间中的表示。"),
            ("word2vec", "一种学习 word embedding 的方法族，包括 skip-gram 与 CBOW。"),
            ("skip-gram", "用中心词邻域构造上下文的训练思路。"),
            ("contextual embedding", "上下文嵌入：同一词在不同上下文中拥有不同向量，如 BERT。"),
        ],
        "methods": [
            "PCA / SVD：对 term-document matrix 做分解，选择最大主成分，将高维稀疏向量投影到低维空间。",
            "skip-gram 直觉：对中心词 w 与上下文词 c 的共现关系建模；通过正例和负例训练，使相似上下文中的词靠近。",
            "类比推理：x = vector(tree) - vector(apple) + vector(grape)，再用 cosine similarity 找与 x 最接近的词。",
            "基础文档向量：vector(document) = sum_i vector(word_i)。这是简单基线，不保留复杂词序信息。",
        ],
        "review": [
            "区分高维稀疏 BOW / TF-IDF 与低维稠密 embedding 的优缺点。",
            "会解释 PCA 的目的、SVD 的角色以及二维可视化的含义。",
            "掌握 distributional semantics、skip-gram、word2vec、contextual embedding 和 BERT 的关系。",
        ],
    },
    "Lecture 18 -- Intro to Transformers.pdf": {
        "title": "Lecture 18: Transformer 入门 / Intro to Transformers",
        "overview": (
            "本讲逐层搭建 Transformer architecture。核心构件包括残差连接（residual connection）、"
            "注意力（attention）、位置编码（positional encoding）、自注意力（self-attention）、"
            "多头注意力（multi-head attention）、前馈网络（feed-forward network）和层归一化"
            "（layer normalization）。"
        ),
        "points": [
            "残差连接（residual connection）跳过一层或多层，为前向传播提供替代路径，也为反向传播提供更短路径，有助于缓解 vanishing gradient。",
            "注意力（attention）学习输入序列上不同 token 的权重分布，使模型聚焦于当前任务相关的上下文。",
            "Transformer 回到固定长度序列：输入长度没有理论上限，但越长的序列会增加模型复杂度和训练成本。",
            "attention 本身不编码顺序，因此必须加入位置编码（positional encoding）。",
            "正弦/余弦位置编码（sinusoidal positional encoding）是静态连续表示；BERT 与 GPT 可使用 learned positional embedding。",
            "词嵌入与位置嵌入通常通过逐元素相加聚合，再送入网络。",
            "encoder 包含 embedding、positional encoding、multi-head attention、residual connection、layer normalization 和 feed-forward network。",
            "前馈网络常先扩张再收缩隐藏维度，例如 512 -> 2048 -> 512。",
            "自注意力（self-attention）通过线性变换、缩放与 softmax 得到 token 之间的依赖概率。",
            "多头注意力（multi-head attention）并行学习多种依赖关系；encoder layer 可堆叠多层。",
        ],
        "terms": [
            ("Transformer", "基于 attention 的序列建模架构，不依赖循环网络处理 token 依赖。"),
            ("residual connection", "残差连接：跳过部分层并聚合向量，改善深层网络训练。"),
            ("attention", "注意力：学习上下文 token 的重要性权重。"),
            ("positional encoding", "位置编码：为序列中的 token 加入顺序信息。"),
            ("self-attention", "自注意力：在同一输入序列内部建模 token 间依赖。"),
            ("multi-head attention", "多头注意力：并行使用多个 attention head 学习不同关系。"),
            ("layer normalization", "层归一化：稳定每层表示的数值分布。"),
            ("BERT", "基于 Transformer 的 contextual embedding 模型。"),
        ],
        "methods": [
            "残差聚合常写为 y = F(x) + x；也可用 concatenation，但会改变下一层宽度。",
            "位置输入：input_embedding(position_i) = word_embedding(token_i) + positional_embedding(i)。",
            "scaled dot-product self-attention：Attention(Q,K,V) = softmax(QK^T / sqrt(d_k))V。",
            "multi-head attention：head_i = Attention(QW_i^Q, KW_i^K, VW_i^V)，再将多个 head 拼接并线性映射。",
            "encoder block：multi-head attention -> residual + layer normalization -> feed-forward network -> residual + layer normalization。",
        ],
        "review": [
            "会按数据流解释 Transformer encoder 的每个模块及其作用。",
            "能说明为什么 attention 仍需要 positional encoding，并比较 sinusoidal 与 learned positional embedding。",
            "记住 self-attention 公式中 Q、K、V、softmax 与 sqrt(d_k) 缩放的作用。",
            "理解 multi-head attention 为什么可以捕捉多种 token 依赖关系。",
        ],
    },
}

def make_note(title, overview, points, terms, methods, review):
    return {
        "title": title,
        "overview": overview,
        "points": points,
        "terms": terms,
        "methods": methods,
        "review": review,
    }


LECTURE_NOTES.update({
    "Lecture 08 -- Intro to ML, Document Clustering.pdf": make_note(
        "Lecture 08: 机器学习与文档聚类 / ML and Document Clustering",
        "本讲从机器学习（Machine Learning）的基本工作流进入无监督学习（unsupervised learning），重点介绍 K-Means 与层次聚类（hierarchical clustering）。",
        [
            "监督学习需要 training、validation、testing 三个集合；validation 用于调参，testing 只在最后客观评估一次。",
            "聚类（clustering）把相似对象放入同一组，可用于数据探索、预处理、索引、压缩和缺少标签时的分类。",
            "K-Means 预先指定 k 个 cluster，交替执行“分配到最近 centroid”和“重算均值”，直到收敛。",
            "K-Means 对初始化、不同密度、非球形 cluster 和 outlier 敏感；K-Means++ 与多次随机初始化可改善结果。",
            "层次聚类可采用 divisive 或 agglomerative 路线，并用 dendrogram 表示层级关系。",
        ],
        [
            ("unsupervised learning", "无监督学习：在没有标签的数据中寻找结构。"),
            ("K-Means", "K 均值：最小化 cluster 内平方距离的聚类方法。"),
            ("centroid", "质心：一个 cluster 中向量的均值。"),
            ("hierarchical clustering", "层次聚类：逐步拆分或合并 cluster 的方法。"),
            ("dendrogram", "树状图：展示层次聚类合并过程与切分阈值。"),
        ],
        [
            "K-Means 目标：min sum_j sum_(x in X_j) ||x - c_j||_2^2。",
            "Lloyd's algorithm：初始化 k 个 centroid -> 分配样本 -> 更新 centroid -> 重复。",
            "agglomerative clustering 可使用 single-link、complete-link 或 average-link 距离。",
        ],
        ["区分 train / validation / test 的用途。", "掌握 Lloyd's algorithm 与 elbow method。", "会说明 K-Means 和 hierarchical clustering 的差异。"],
    ),
    "Lecture 09 --Prelude to Supervised Learning.pdf": make_note(
        "Lecture 09: 监督学习前奏 / Regression and Gradient Descent",
        "本讲介绍 linear regression、logistic regression 与 gradient descent，为后续神经网络训练建立数学基础，并将这些概念连接到文本分类。",
        [
            "linear regression 用参数逼近输入 X 与响应变量 Y 的线性关系，常以 mean squared error (MSE) 作为 cost function。",
            "logistic regression 使用 sigmoid 将线性输出映射到 [0,1]，可解释为二分类概率，并用 threshold 做决策。",
            "gradient descent 在没有解析解时迭代调整参数，使 loss 沿负梯度方向下降。",
            "文本分类输入是 document d 与固定标签集合 C，模型通常输出每个 class 的概率。",
            "监督分类器需要 feature representation、分类函数以及从标注数据学习参数的算法。",
        ],
        [
            ("linear regression", "线性回归：用线性函数拟合连续响应变量。"),
            ("logistic regression", "逻辑回归：用 sigmoid 输出分类概率。"),
            ("gradient descent", "梯度下降：沿负梯度迭代搜索低 loss 参数。"),
            ("learning rate", "学习率：每次参数更新的步长。"),
            ("cost / loss function", "损失函数：衡量预测与真实值差异。"),
        ],
        [
            "线性模型：y_hat = w · x + b。",
            "MSE：J = (1/N) * sum_i (y_i - y_hat_i)^2。",
            "sigmoid：sigma(z) = 1 / (1 + exp(-z))。",
            "梯度下降：theta_(t+1) = theta_t - lambda * grad J(theta_t)。",
        ],
        ["理解 linear regression 与 logistic regression 的用途区别。", "会解释 sigmoid、threshold、MSE 和 learning rate。", "掌握 gradient descent 的更新方向。"],
    ),
    "Lecture 10 -- Supervised ML.pdf": make_note(
        "Lecture 10: 监督学习与神经网络 / Supervised ML and Neural Networks",
        "本讲将 logistic regression 扩展为人工神经网络（artificial neural network, ANN），讨论 neuron、activation、FFNN、softmax、one-hot label 与 cross-entropy loss。",
        [
            "neuron 对输入做加权和、加 bias，再应用非线性 activation function；常见函数包括 sigmoid、tanh 和 ReLU。",
            "全连接前馈网络（feed-forward neural network, FFNN）由多层 neuron 组成，层参数可写为矩阵。",
            "分类任务可分 binary、multiclass 与 multilabel；multiclass 输出常用 softmax 转成概率分布。",
            "训练数据写作 feature matrix X 与 label vector Y；类别标签常转换为 one-hot vector。",
            "分类器常以 cross-entropy loss 衡量预测分布与真实分布的差异。",
        ],
        [
            ("activation function", "激活函数：为神经网络加入非线性表达能力。"),
            ("ReLU", "修正线性单元：输出 max(z, 0)。"),
            ("softmax", "将多个 logit 归一化为类别概率分布。"),
            ("one-hot encoding", "独热编码：仅真实类别位置为 1 的标签向量。"),
            ("cross-entropy loss", "交叉熵损失：比较预测分布与真实标签分布。"),
        ],
        [
            "neuron：a = f(w · x + b)。",
            "ReLU：f(z) = max(z, 0)。",
            "softmax(z_i) = exp(z_i) / sum_j exp(z_j)。",
            "cross entropy：L = -sum_i y_i log(p_i)。",
        ],
        ["会从 logistic regression 解释 neuron。", "区分 binary、multiclass、multilabel。", "记住 softmax 与 cross-entropy 的配合。"],
    ),
    "Lecture 11 -- Deep Learning -- Generalization and Tuning.pdf": make_note(
        "Lecture 11: 泛化、过拟合与调优 / Generalization and Tuning",
        "本讲讨论分类器评估与神经网络的泛化（generalization）。重点包括 confusion matrix、类别不平衡、overfitting、underfitting、early stopping、regularization 与 dropout。",
        [
            "二分类预测可分 TP、TN、FP、FN；不同业务场景对错误类型的代价不同。",
            "accuracy 在类别不平衡时可能误导，评估时要结合 baseline、precision、recall 与 F1。",
            "overfitting 是模型记住训练集细节却无法推广到未见数据；表现为 training 改善而 validation 停滞或恶化。",
            "噪声、错误标签、稀有特征和伪相关都会促进 overfitting。",
            "改善泛化可从更多数据、data augmentation、architecture、early stopping、weight regularization 与 dropout 入手。",
        ],
        [
            ("generalization", "泛化：模型在未见数据上保持性能的能力。"),
            ("overfitting", "过拟合：模型记忆训练细节，验证性能下降。"),
            ("confusion matrix", "混淆矩阵：统计 TP、TN、FP、FN 或多类别对应关系。"),
            ("regularization", "正则化：限制模型复杂度以改善泛化。"),
            ("dropout", "随机失活：训练时随机屏蔽部分 neuron。"),
        ],
        [
            "accuracy = (TP + TN) / (TP + TN + FP + FN)。",
            "precision = TP / (TP + FP)；recall = TP / (TP + FN)。",
            "F1 = 2 * precision * recall / (precision + recall)。",
        ],
        ["会根据任务解释 FP 与 FN 的代价。", "能识别 learning curve 中的 overfitting。", "比较 early stopping、regularization、dropout 与 data augmentation。"],
    ),
    "Lecture 12 -- Sequence Data and Recurrent NNs.pdf": make_note(
        "Lecture 12: 序列数据与循环神经网络 / RNN, GRU and LSTM",
        "本讲解释为什么文本需要序列模型，并介绍 recurrent neural network (RNN)、GRU、LSTM、encoder-decoder、BRNN 与 vanishing gradient。",
        [
            "BOW 忽略顺序，n-gram 只能建模短上下文；自然语言则包含长距离依赖。",
            "RNN 将前一时间步 activation 反馈到当前输入，用共享权重逐 token 处理序列。",
            "GRU 使用 recall gate 与 update gate 管理 memory；LSTM 进一步分离 memory cycle，并使用多个 gate。",
            "RNN 可组织成 sequence-to-sequence、sequence-to-vector、vector-to-sequence 和 encoder-decoder。",
            "双向 RNN（BRNN）合并正向和反向上下文，但深层循环网络仍面临训练慢与 vanishing gradient。",
        ],
        [
            ("RNN", "循环神经网络：通过反馈 activation 处理序列。"),
            ("GRU", "门控循环单元：用 gate 管理记忆读写。"),
            ("LSTM", "长短期记忆网络：具有独立 memory cycle 的 RNN。"),
            ("encoder-decoder", "编码器-解码器：先压缩输入序列，再生成输出序列。"),
            ("vanishing gradient", "梯度消失：长路径反向传播时梯度逐渐衰减。"),
        ],
        [
            "基本 RNN：h_t = tanh(W_x x_t + W_h h_(t-1) + b)。",
            "sequence-to-vector 适合分类；vector-to-sequence 适合生成；encoder-decoder 适合翻译。",
        ],
        ["理解 RNN 展开图与共享权重。", "比较 RNN、GRU、LSTM。", "会为 POS tagging、分类、captioning、translation 选择架构。"],
    ),
    "Lecture 13 -- Applications of RNNs.pdf": make_note(
        "Lecture 13: BRNN、序列标注与生成 / Applications of RNNs",
        "本讲用 BRNN 处理序列标注，重点介绍 POS tagging、HMM、Viterbi、NER、BIO tagging，以及生成语言模型中的 beam search。",
        [
            "BRNN 同时运行正向与反向 RNN，并拼接 activation，使每个位置获得两侧上下文。",
            "POS tagging 为每个 token 分配词性标签；常见方法包括 HMM、CRF、RNN、Transformer 和微调 BERT。",
            "HMM 将标签视为 hidden state，将词视为 observation；Viterbi 用 dynamic programming 求最可能标签路径。",
            "NER 识别人名、地点、组织等实体 span；BIO tagging 将实体边界转换为逐 token 标签。",
            "生成模型逐步预测 next token；beam search 保留多个候选序列，平衡效率与全局质量。",
        ],
        [
            ("POS tagging", "词性标注：按上下文标记每个 token 的语法类别。"),
            ("Hidden Markov Model (HMM)", "隐马尔可夫模型：具有隐藏状态与发射概率的序列模型。"),
            ("Viterbi algorithm", "维特比算法：动态规划求 HMM 最可能路径。"),
            ("BIO tagging", "实体边界标注：B 表示开始，I 表示内部，O 表示实体外。"),
            ("beam search", "束搜索：每步保留若干高分候选序列。"),
        ],
        [
            "HMM 参数：初始概率 pi、转移概率 A、发射概率 B。",
            "Viterbi：delta_t(s) = max_(s') delta_(t-1)(s') * A(s',s) * B(s,o_t)。",
            "序列 loss 通常取 token 平均 log loss，避免长序列天然受到更大惩罚。",
        ],
        ["会解释 HMM 的 hidden state、observation、transition、emission。", "能为实体短语写 BIO 标签。", "理解 greedy decoding 与 beam search 的差别。"],
    ),
    "Lecture 14 -- Generative Models.pdf": make_note(
        "Lecture 14: 生成模型与 word2vec / Generative Models and Embeddings",
        "本讲连接 RNN 生成、beam search 与 word embedding，重点比较 word2vec 的 skip-gram 和 CBOW。",
        [
            "RNN language model 从 <s> 开始逐步生成 token，直到 </s>；训练数据可由句子滑动窗口构造。",
            "beam search 设置 beam width、expansion factor 与 goodness metric，保留有限个高分候选。",
            "embedding 用短而稠密的向量替代高维稀疏 TF / TF-IDF 表示。",
            "word2vec skip-gram 根据中心词预测上下文；CBOW 根据上下文预测中心词。",
            "skip-gram 对小语料和稀有词较有优势；CBOW 更快，对常见词表现较好。",
        ],
        [
            ("generative model", "生成模型：学习序列分布并产生新样本。"),
            ("beam width", "束宽：beam search 每轮保留的候选数。"),
            ("skip-gram", "根据目标词预测附近上下文。"),
            ("continuous bag-of-words (CBOW)", "根据上下文预测目标词。"),
            ("negative sampling", "负采样：用非共现 pair 降低训练成本。"),
        ],
        [
            "beam search：扩展候选 -> 用 perplexity 或 log probability 排序 -> 保留前 N 个。",
            "skip-gram 的 embedding 是隐藏层中与目标词对应的 weight vector。",
        ],
        ["比较 greedy sampling 与 beam search。", "区分 skip-gram 和 CBOW 的输入输出。", "理解 embedding 为何比稀疏向量更适合迁移。"],
    ),
    "Lecture 15 -- MachineTranslation.pdf": make_note(
        "Lecture 15: 词嵌入、机器翻译与注意力 / Machine Translation and Attention",
        "本讲补充 GloVe、fastText、Doc2Vec，并从机器翻译引出 attention mechanism：模型应为每个输出 token 动态关注输入序列中的相关位置。",
        [
            "GloVe 从 co-occurrence matrix 出发，经频率调整、gradient descent 与降维得到 embedding。",
            "fastText 在 word2vec 思路上加入 character n-gram，因此能利用词内部结构并缓解 unknown word 问题。",
            "sentence / document embedding 可使用平均 word embedding、加权平均或 Doc2Vec。",
            "传统 encoder-decoder 的单一 context vector 容量有限，远距离信息容易衰减。",
            "attention 为每个输出 token 计算输入位置上的概率权重，再聚合成动态 context vector。",
        ],
        [
            ("GloVe", "Global Vectors：基于全局共现统计学习 embedding。"),
            ("fastText", "使用 character n-gram 的词嵌入方法。"),
            ("Doc2Vec", "将 paragraph ID 加入训练的文档嵌入方法。"),
            ("attention", "注意力：对输入位置分配权重并聚合相关信息。"),
            ("context vector", "上下文向量：encoder 信息传给 decoder 的表示。"),
        ],
        [
            "attention 权重：alpha_i = softmax(e_i)。",
            "动态上下文：c = sum_i alpha_i * h_i，其中 h_i 是 encoder activation。",
        ],
        ["会比较 word2vec、GloVe、fastText。", "理解 attention 解决固定 context vector 瓶颈的方式。", "能解释 attention matrix 的行列含义。"],
    ),
    "Lecture 16 -- NeuralNetwork.Advanced.pdf": make_note(
        "Lecture 16: 神经网络调优与高级特性 / NN Tuning",
        "本讲给出训练神经网络的工程实践：模型保存、cross-validation、optimizer、regularization、normalization、early stopping 与 batch size。",
        [
            "避免重复计算：保存预处理结果与 best model，测试时加载 validation 表现最优的 checkpoint。",
            "K-fold cross-validation 循环选择 fold 做 validation，适合小数据集，但计算成本较高。",
            "optimizer 决定 gradient 更新策略；SGD 简单，Adagrad、RMSProp、Adam 使用历史梯度自适应学习率。",
            "L1、L2 与 dropout 通过限制复杂度缓解 overfitting；PyTorch 中 L2 常对应 weight_decay。",
            "layer normalization 与 batch normalization 可稳定梯度、加快收敛；early stopping 避免继续训练导致过拟合。",
        ],
        [
            ("cross-validation", "交叉验证：轮流使用不同 fold 做验证。"),
            ("optimizer", "优化器：根据梯度更新参数的规则。"),
            ("Adam", "结合一阶与二阶动量的自适应 optimizer。"),
            ("L2 regularization", "对较大参数加平方惩罚。"),
            ("layer normalization", "按单个样本的层输出做归一化。"),
        ],
        [
            "L2 penalty：L_total = L + lambda * sum_i w_i^2。",
            "normalization：x_hat = (x - mu) / sigma，再学习 scale 与 shift。",
        ],
        ["比较 static validation 与 cross-validation。", "理解 SGD、Adam、learning-rate schedule。", "掌握 dropout、weight_decay、early stopping 的使用目的。"],
    ),
    "Lecture 17 -- NeuralNetwork.Advanced.Features 2.pdf": make_note(
        "Lecture 17: CNN、Embedding Layer 与网络几何 / Advanced NN Features",
        "本讲介绍 convolution、pooling、embedding layer，并讨论 NLP 网络应多深、多宽。核心观点是模型容量要与任务匹配，而不是盲目增加 neuron。",
        [
            "CNN 用 kernel 在局部区域滑动计算 dot product，pooling 再做局部降维；图像任务中常交替堆叠。",
            "一维 CNN 可识别文本中的局部 n-gram pattern，也可用于 DNA 等线性序列。",
            "embedding layer 是 index 到 dense vector 的 lookup table，可随机初始化、加载 pretrained embedding 或继续 fine-tune。",
            "文本有序列依赖，FFNN 与 RNN 对长距离依赖有限；Transformer 改变了深层 NLP 网络的能力边界。",
            "增加宽度会提高判别能力，但也增加训练时间与 overfitting 风险。",
        ],
        [
            ("convolution", "卷积：用共享 kernel 扫描局部区域提取模式。"),
            ("pooling", "池化：对局部区域取最大值或平均值进行降维。"),
            ("embedding layer", "嵌入层：将 token index 映射到 dense vector。"),
            ("pretrained embedding", "预训练嵌入：从外部语料预先学习的向量。"),
        ],
        [
            "convolution 输出是 kernel 与局部窗口的 dot product。",
            "embedding matrix 形状通常为 |V| x D，其中 D 是 embedding dimension。",
        ],
        ["会解释 CNN 在文本中的 1D n-gram 直觉。", "理解 embedding layer 的三种初始化路线。", "能说明网络过宽、过深的代价。"],
    ),
    "Lecture 19 -- Intro to Transformers and BERT.pdf": make_note(
        "Lecture 19: Transformer 解码器、采样与 BERT / Sampling and BERT",
        "本讲补全 Transformer decoder，并介绍生成采样、transfer learning 与 BERT 预训练。",
        [
            "生成模型的 sampling 策略包括原始分布采样、greedy、temperature、top-k 和 top-p。",
            "Transformer decoder 使用 masked multi-head attention 模拟从左到右生成，并接收 encoder context。",
            "transfer learning 包括 feature-based 与 fine-tuning：前者冻结预训练模型，后者继续调整其参数。",
            "GPT 使用 causal 左上下文；BERT 使用双向 encoder 表示。",
            "BERT 通过 masked language model (MLM) 与 next sentence prediction (NSP) 预训练。",
        ],
        [
            ("temperature", "温度：调节概率分布尖锐或平坦程度。"),
            ("top-k sampling", "仅在概率最高的 k 个 token 中采样。"),
            ("top-p sampling", "在累计概率达到 p 的最小候选集合中采样。"),
            ("fine-tuning", "微调：继续训练预训练模型以适配下游任务。"),
            ("masked language model (MLM)", "遮盖语言模型：根据双向上下文预测被遮盖 token。"),
        ],
        [
            "temperature softmax：p_i = exp(z_i / T) / sum_j exp(z_j / T)。",
            "BERT MLM：抽取 15% token，其中 80% 换成 [MASK]，10% 随机替换，10% 保持不变。",
        ],
        ["比较 greedy、temperature、top-k、top-p。", "区分 encoder 与 masked decoder。", "掌握 BERT 的 MLM、NSP、[CLS]、[SEP]。"],
    ),
    "Lecture 20 -- Bert and Friends NLP Tasks.pdf": make_note(
        "Lecture 20: Transformer 家族与 BERT 任务 / BERT and Friends",
        "本讲按 encoder-only、decoder-only 与 encoder-decoder 梳理 Transformer 家族，并讨论微调风险、参数高效适配、DistilBERT 和 BERT 的 NLP 应用。",
        [
            "冻结 pretrained model 可把它当 feature extractor；完全 fine-tune 则更灵活，但成本高且可能 catastrophic forgetting。",
            "Transformer 家族可分 encoder-based（BERT 等）、decoder-based（GPT 等）与 encoder-decoder（T5、BART 等）。",
            "DistilBERT 用 knowledge distillation 让较小 student 模型模仿 BERT teacher，减少参数并加快推理。",
            "BERT 可用于 sentence classification、sentence-pair relationship、sequence labeling 和 NER。",
            "NLP 应用还包括 sentiment analysis、IR、knowledge graph、summarization、QA、reasoning 和多模态转换。",
        ],
        [
            ("catastrophic forgetting", "灾难性遗忘：微调破坏预训练阶段获得的能力。"),
            ("parameter-efficient fine-tuning (PEFT)", "参数高效微调：只更新较少附加参数。"),
            ("knowledge distillation", "知识蒸馏：student 模型学习 teacher 模型行为。"),
            ("DistilBERT", "蒸馏后的轻量 BERT。"),
        ],
        [
            "BERT base：12 layers、12 heads、768 hidden units、约 110M parameters。",
            "BERT large：24 layers、16 heads、1024 hidden units、约 340M parameters。",
        ],
        ["按架构分类 BERT、GPT、T5。", "理解 frozen feature extractor 与 fine-tuning 的权衡。", "说明 knowledge distillation 的目标。"],
    ),
    "Lecture 21 -- GPT and T5.pdf": make_note(
        "Lecture 21: GPT、T5 与 Transformer 家族 / GPT and T5",
        "本讲比较 decoder-only GPT 与 encoder-decoder T5，讨论模型规模、zero-shot learning、emergent properties 与 scaling law。",
        [
            "GPT 是 decoder-only causal Transformer，通过 masked attention 做 autoregressive next-token prediction。",
            "GPT-2 在大规模数据训练后表现出 zero-shot learning，可在未显式提供任务样例时执行问答、摘要和翻译。",
            "大模型可能出现 emergent properties，也可能生成虚构事实，因此能力与风险同时增长。",
            "scaling law 描述 loss 与 parameter、dataset size、compute 之间较平滑的幂律关系。",
            "T5 保留 encoder-decoder 架构，把多种任务统一为 text-to-text，并以 denoising 目标训练。",
        ],
        [
            ("autoregressive generation", "自回归生成：根据已有左上下文逐步预测下一 token。"),
            ("zero-shot learning", "零样本学习：没有特定任务示例也能执行任务。"),
            ("emergent properties", "涌现能力：规模增长后出现未显式设计的行为。"),
            ("scaling law", "缩放规律：性能与参数、数据、算力之间的规律关系。"),
            ("T5", "Text-to-Text Transfer Transformer：统一文本输入输出的模型。"),
        ],
        [
            "causal LM：P(W) = product_t P(w_t | w_1,...,w_(t-1))。",
            "T5 denoising：破坏输入文本片段，再训练模型恢复原文。",
        ],
        ["比较 BERT、GPT、T5 的架构与预训练目标。", "理解 zero-shot 与 emergent properties。", "会解释 scaling law 的三个主要变量。"],
    ),
    "Lecture 22 -- Audio Processing and Deep Learning.pdf": make_note(
        "Lecture 22: 音频处理与深度学习 / Audio Processing",
        "本讲从 GPT 扩展到音频处理基础：声波、采样、Fourier transform、spectrum、spectrogram 与 Mel scale，为 ASR 建立信号表示。",
        [
            "声音是空气压力随时间传播的纵波；纯正弦波可由 wavelength、period、amplitude 与 frequency 描述。",
            "数字音频以固定 sample rate 对连续模拟信号采样，得到离散 amplitude 序列。",
            "Fourier analysis 将 time domain 信号分解为 frequency domain 中的正弦分量；FFT 可高效实现。",
            "spectrogram 记录频率成分随时间变化，可用 heat map 显示 amplitude。",
            "人类对音高与响度的感知接近 log scale，因此语音处理中常使用 log Mel spectrogram。",
        ],
        [
            ("frequency", "频率：每秒周期数，单位 Hertz (Hz)。"),
            ("Fourier transform", "傅里叶变换：在 time domain 与 frequency domain 间转换。"),
            ("spectrum", "频谱：展示频率分量与 amplitude。"),
            ("spectrogram", "声谱图：展示频率、amplitude 随时间变化。"),
            ("Mel scale", "梅尔刻度：近似人类音高感知的非线性刻度。"),
        ],
        [
            "frequency 与 period：f = 1 / p。",
            "波长关系：v = f * lambda。",
            "强度近似：I = c * A^2。",
            "FFT 复杂度约为 O(N log N)。",
        ],
        ["区分 time domain、frequency domain、spectrum、spectrogram。", "记住 f、p、lambda、A 的意义。", "说明为何 ASR 常使用 log Mel spectrogram。"],
    ),
    "Lecture 23 -- ASR.pdf": make_note(
        "Lecture 23: 自动语音识别 / Automatic Speech Recognition",
        "本讲介绍自动语音识别（Automatic Speech Recognition, ASR）：从 phoneme、IPA、formant、prosody 到 HMM/Viterbi 与 Transformer-based ASR。",
        [
            "phoneme 是特定语言中能够区别意义的最小声音单位；IPA 是跨语言语音符号系统，ARPAbet 是 ASCII 版本。",
            "vowel 的频谱具有 characteristic formant；consonant 与 semi-vowel 更依赖随时间变化的 spectrogram。",
            "连续语音识别要处理 phoneme、word、句法、语义、prosody、speaker、accent 和 self-correction。",
            "传统 ASR 将 log Mel spectrogram 转成 MFCC 等 feature sequence，再交给 HMM/Viterbi，并结合 n-gram language model 解码。",
            "现代 ASR 将任务视为 sequence-to-sequence，常用 RNN 或 Transformer；Whisper 是 Transformer-based ASR 示例。",
        ],
        [
            ("ASR", "自动语音识别：将语音信号转换为文本。"),
            ("phoneme", "音位：在一种语言中区别意义的最小声音单位。"),
            ("formant", "共振峰：由声道形状造成的频谱特征峰。"),
            ("prosody", "韵律：语音中的重音、时长、音高变化等特征。"),
            ("MFCC", "Mel-frequency cepstral coefficients：常用语音特征。"),
            ("word error rate (WER)", "词错误率：ASR 输出与参考文本差异指标。"),
        ],
        [
            "HMM ASR 流程：audio -> log Mel spectrogram -> feature vectors / MFCC -> HMM observation sequence -> Viterbi decoding + n-gram LM。",
            "WER = (substitutions + deletions + insertions) / reference_words。",
        ],
        ["掌握 phoneme、IPA、formant、prosody。", "会描述传统 HMM ASR 管线。", "理解 Transformer-based ASR 与传统方法的关系。"],
    ),
    "Lecture 24 -- Last Lecture.pdf": make_note(
        "Lecture 24: NLP 的未来 / The Future of NLP",
        "最后一讲回顾课程脉络，并讨论大模型未来、AI 风险、Turing Test、LaMDA、人工通用智能（AGI）、Chinese Room 与 emergent properties。",
        [
            "NLP 的进展由硬件、软件、算法与数据共同推动；单纯扩大 parameter count 可能逐渐遇到边际收益下降。",
            "未来方向包括更好的数据筛选、训练算法、transfer learning、事实核查、bias checking、prompt engineering 与多模态能力。",
            "AI 风险讨论不仅关乎性能，还涉及可控性、依赖、错误信息、偏见与社会影响。",
            "Turing Test 衡量可观察对话行为，不等同于证明机器具有意识或理解。",
            "Chinese Room 思想实验追问符号处理是否构成真正理解；emergent properties 则提醒我们复杂系统可能出现未预期行为。",
        ],
        [
            ("artificial general intelligence (AGI)", "人工通用智能：跨任务表现出广泛适应能力的 AI。"),
            ("bias checking", "偏差检查：识别模型输出中的系统性偏见。"),
            ("Chinese Room", "中文房间：质疑符号操作是否等于语义理解的思想实验。"),
            ("emergent properties", "涌现性质：系统整体出现组成部分未直接体现的行为。"),
        ],
        [
            "分析框架：能力提升 -> 应用价值 -> 错误与偏见 -> 可控性 -> 社会治理。",
            "讨论意识时应区分 behavioral evidence、内部机制与哲学定义。",
        ],
        ["能从全课程串联 preprocessing、vector model、RNN、attention、Transformer、LLM、ASR。", "会批判性讨论 Turing Test 的边界。", "能列出大模型未来机会与风险。"],
    ),
})

NARRATIVE_INTROS = {
    "Lecture 02 -- What is NLP.pdf": "我们先从一个看似简单的问题出发：计算机能否真正处理人类语言？语言不是整齐的表格，而是带有歧义、上下文和意图的表达。于是 NLP 的第一步不是立刻选择模型，而是先看清需要解决的任务版图：理解、检索、分类、生成、问答和多模态转换。",
    "Lecture 03 -- Basic Notions and Text Normalization.pdf": "模型无法直接理解一份杂乱的网页、聊天记录或扫描文本。缩写、金额、URL、表情和不同语言的分词习惯都会让简单切割失败。解决过程从清洗开始：把原始文本逐层整理成 character、token、sentence、document 和 corpus。",
    "Lecture 05 -- Language Models, Smoothing, Perplexity.pdf": "当我们想让机器预测下一个词时，最直接的方法是统计过去出现过的局部组合。但训练语料只是语言世界的一小部分：句子越长，概率乘积越小；测试集还会出现训练时没见过的组合。本讲沿着“如何评价模型”和“如何处理未知事件”两个问题展开。",
    "Lecture 06 -- Introduction to Term Vector Models.pdf": "n-gram 的零概率问题逼迫我们重新分配概率质量，这带来了 smoothing。接着出现另一个问题：如果目标不再是生成句子，而是比较文档，我们需要把文本放进一个可以计算距离的空间。BOW、TF-IDF 和 cosine similarity 就是这条路线的基础工具。",
    "Lecture 07 -- Vector Models, Word and Text Embeddings.pdf": "词频向量能工作，但它们太长、太稀疏，也不真正理解“car”和“automobile”相近。解决思路是让词根据上下文学习一个短而稠密的位置：相似语境中的词靠近，不同语境中的词远离。这一步把统计计数推进到 distributional semantics 和 word embedding。",
    "Lecture 08 -- Intro to ML, Document Clustering.pdf": "当文档数量很大而人工标签稀缺时，我们仍然希望知道数据中是否存在自然分组。无监督学习先不问“正确答案是什么”，而是问“哪些对象彼此相似”。K-Means 与 hierarchical clustering 提供了两种组织数据的方式。",
    "Lecture 09 --Prelude to Supervised Learning.pdf": "如果已经有标签，问题就从“发现结构”变成“学习决策边界”。为了理解分类器如何学习，本讲先从 regression 开始：定义误差，再让参数沿着降低误差的方向更新。gradient descent 是后续神经网络训练的共同基础。",
    "Lecture 10 -- Supervised ML.pdf": "监督学习真正困难的地方不是记住训练样本，而是在未见数据上仍然做出正确判断。我们需要选择特征、损失函数和分类器，并控制 overfitting。Naive Bayes、logistic regression、softmax 和 evaluation metrics 构成了这套基本工具箱。",
    "Lecture 11 -- Deep Learning -- Generalization and Tuning.pdf": "神经网络层数增加后，表达能力变强，但训练也更容易失控：模型可能过拟合、梯度不稳定，或者在错误的超参数上浪费计算。解决方法不是盲目加层，而是系统地调节 learning rate、batch、regularization、dropout 和 validation 策略。",
    "Lecture 12 -- Sequence Data and Recurrent NNs.pdf": "普通前馈网络把输入看作固定大小的向量，但语言天然有顺序：前面的词会影响后面的理解。RNN 引入隐藏状态，把过去的信息带到当前时间步。这样，模型第一次拥有了处理可变长度 sequence 的基础机制。",
    "Lecture 13 -- Applications of RNNs.pdf": "有了 RNN 之后，下一步是把它用于真实任务：序列分类、语言模型、生成和 tagging。但长距离依赖会让普通 RNN 的记忆衰减。LSTM、GRU 和双向结构尝试保留更有价值的信息，并让模型看到更完整的上下文。",
    "Lecture 14 -- Generative Models.pdf": "生成模型不仅要判断一句话好不好，还要逐步产生新的 token。问题在于：每一步都要在概率分布中选择下一个 token，过于贪心会单调，过于随机会失控。temperature、sampling 和 beam search 是在稳定性与多样性之间做权衡的工具。",
    "Lecture 15 -- MachineTranslation.pdf": "机器翻译表面上是把一种语言替换成另一种语言，实际却要处理词序、歧义和长距离依赖。早期统计方法逐步让位于 encoder-decoder 神经网络：encoder 压缩源句，decoder 根据上下文逐步生成目标句，attention 再帮助模型避免遗忘。",
    "Lecture 16 -- NeuralNetwork.Advanced.pdf": "深层网络训练时，梯度可能逐层变小或变大，导致学习停滞或数值不稳定。为了让优化过程可控，本讲加入更成熟的训练机制：更合适的 initialization、normalization、optimizer 和梯度处理方法。",
    "Lecture 17 -- NeuralNetwork.Advanced.Features 2.pdf": "即使训练稳定，序列模型仍然面临信息瓶颈：一个固定长度向量很难完整保存长句。attention 改变了问题的解法：decoder 不必只依赖单个摘要，而可以在每一步回头查看最相关的输入位置。",
    "Lecture 18 -- Intro to Transformers.pdf": "attention 已经缓解了序列瓶颈，但 RNN 仍然必须按时间步顺序计算。Transformer 的关键尝试是去掉 recurrence，让 token 之间直接建立依赖。为此还要补回顺序信息，并用 residual connection 保证深层训练。",
    "Lecture 19 -- Intro to Transformers and BERT.pdf": "Transformer encoder 可以同时理解左右上下文，这为通用语言表示提供了新路径。BERT 的思路是先在海量文本上做预训练，再用少量标注数据 fine-tune 到具体任务。问题从“每个任务训练一个模型”转向“先学语言，再适配任务”。",
    "Lecture 20 -- Bert and Friends NLP Tasks.pdf": "预训练模型的价值最终要落到任务上。分类、NER、问答和文本推断需要不同输出形式，但不必从头训练语言理解能力。本讲关注如何在共享 Transformer 表示上接入 task-specific head，并用 fine-tuning 完成迁移。",
    "Lecture 21 -- GPT and T5.pdf": "BERT 擅长理解上下文，而生成任务需要模型持续预测下一个 token。GPT 采用 decoder-only 自回归路线，T5 则把各种任务统一成 text-to-text。随着模型、数据和算力扩大，语言模型开始展现更强的 few-shot 与 zero-shot 能力。",
    "Lecture 22 -- Audio Processing and Deep Learning.pdf": "语音是 NLP 的另一种入口，但音频不是文本，而是随时间变化的压力波。要让模型使用声音，必须先把波形转成结构化特征。Fourier transform、spectrogram 和 Mel scale 把物理信号转换为更接近人类听觉的表示。",
    "Lecture 23 -- ASR.pdf": "自动语音识别（ASR）要把连续声音还原为文本。难点不仅是识别 phoneme，还包括口音、语速、停顿、韵律和语言上下文。传统 HMM 管线把声学特征与 language model 组合起来，现代 Transformer 则用端到端训练进一步简化流程。",
    "Lecture 24 -- Last Lecture.pdf": "课程最后回到最初的问题：当 NLP 系统越来越像是在理解语言，我们应该如何判断它们的能力和边界？未来进展不只取决于参数规模，还取决于数据、推理效率、对齐和社会治理。本讲把技术路线连接到更长期的问题。",
}

REAL_WORLD_APPLICATIONS = {
    "Lecture 02 -- What is NLP.pdf": ["搜索引擎、客服机器人、内容审核、自动摘要、机器翻译和智能写作都属于 NLP 任务版图。", "企业落地时通常不是选择一个万能模型，而是把分类、检索、抽取和生成组合成流程。"],
    "Lecture 03 -- Basic Notions and Text Normalization.pdf": ["搜索索引、日志分析、社交媒体监测和 RAG 数据入库都依赖可靠的 cleaning 与 tokenization。", "输入规范化不足会直接造成召回率下降、重复数据增加和模型评估失真。"],
    "Lecture 05 -- Language Models, Smoothing, Perplexity.pdf": ["perplexity 可用于比较 language model 版本、检查领域适配效果和发现数据分布漂移。", "拼写纠正、输入法联想和早期语音识别系统都使用过 n-gram language model。"],
    "Lecture 06 -- Introduction to Term Vector Models.pdf": ["TF-IDF 与 cosine similarity 仍是搜索、文档去重、相似案例检索和轻量级推荐系统的强基线。", "在数据量较小或需要可解释性时，传统向量模型通常比复杂模型更容易部署和审计。"],
    "Lecture 07 -- Vector Models, Word and Text Embeddings.pdf": ["embedding 支撑语义搜索、推荐系统、聚类、异常检测和 RAG 向量数据库。", "contextual embedding 能提升同义词检索，并根据上下文区分一词多义。"],
    "Lecture 08 -- Intro to ML, Document Clustering.pdf": ["文档聚类可用于新闻主题整理、客户反馈归类、知识库清理和数据标注前的探索。", "聚类结果常用于发现标签体系缺口，而不是直接当作最终业务答案。"],
    "Lecture 09 --Prelude to Supervised Learning.pdf": ["gradient descent 是广告排序、风险预测、文本分类和深度学习训练的共同优化基础。", "logistic regression 因可解释、稳定和部署成本低，仍常用于生产基线。"],
    "Lecture 10 -- Supervised ML.pdf": ["垃圾邮件识别、工单路由、情感分析、欺诈检测和医疗文本分类都是监督学习场景。", "真实系统必须同时关注 precision、recall、F1 与类别不平衡，而不是只看 accuracy。"],
    "Lecture 11 -- Deep Learning -- Generalization and Tuning.pdf": ["超参数调优决定模型能否稳定上线；validation、early stopping 和 dropout 常用于控制过拟合。", "算力预算有限时，系统化实验记录比盲目扩大模型更重要。"],
    "Lecture 12 -- Sequence Data and Recurrent NNs.pdf": ["RNN 曾广泛用于语言模型、时间序列预测、日志异常检测和语音处理。", "理解 hidden state 仍有助于分析现代 sequence model 如何保存和更新上下文。"],
    "Lecture 13 -- Applications of RNNs.pdf": ["LSTM 与 GRU 适合处理传感器数据、金融序列、早期机器翻译和序列标注。", "双向模型适合离线分析，但不适合严格实时、不能查看未来输入的场景。"],
    "Lecture 14 -- Generative Models.pdf": ["temperature、top-k、top-p 和 beam search 直接影响聊天机器人、摘要系统与翻译系统的输出风格。", "生产环境通常需要在多样性、事实性、延迟和安全约束之间调参。"],
    "Lecture 15 -- MachineTranslation.pdf": ["机器翻译用于跨语言客服、本地化、字幕和多语言知识库。", "attention 的思想也延伸到摘要、问答和多模态模型。"],
    "Lecture 16 -- NeuralNetwork.Advanced.pdf": ["normalization、optimizer 和 initialization 决定大型训练任务能否稳定收敛。", "gradient clipping 在语音、序列和生成模型中常用于避免训练爆炸。"],
    "Lecture 17 -- NeuralNetwork.Advanced.Features 2.pdf": ["attention 可用于机器翻译对齐、文档摘要、图像描述和多模态检索。", "attention map 也可作为分析工具，但不能简单等同于完整解释。"],
    "Lecture 18 -- Intro to Transformers.pdf": ["Transformer 已成为搜索排序、文本生成、代码助手、视觉语言模型和语音模型的核心架构。", "self-attention 允许并行训练，但长序列成本高，推动了稀疏 attention 与长上下文优化。"],
    "Lecture 19 -- Intro to Transformers and BERT.pdf": ["BERT 类模型适合分类、NER、语义匹配和抽取式 QA。", "预训练加 fine-tuning 降低了每个任务独立收集大规模标签的成本。"],
    "Lecture 20 -- Bert and Friends NLP Tasks.pdf": ["共享 encoder 加 task head 是企业文本分类、实体识别和问答系统的常见实现方式。", "在标注数据有限时，可先从预训练模型和小规模 fine-tuning 开始。"],
    "Lecture 21 -- GPT and T5.pdf": ["GPT 驱动对话、写作、代码生成和 agent；T5 的 text-to-text 思路适合统一多任务接口。", "模型规模扩大带来能力提升，也增加推理成本、幻觉与治理压力。"],
    "Lecture 22 -- Audio Processing and Deep Learning.pdf": ["Mel spectrogram 是语音识别、说话人识别、音乐分类和声音事件检测的常用输入。", "频域表示能让模型更容易捕捉音高、音色和随时间变化的模式。"],
    "Lecture 23 -- ASR.pdf": ["ASR 用于会议转录、字幕、语音助手、呼叫中心质检和无障碍工具。", "现代系统如 Whisper 使用 Transformer 处理多语言和长音频转录。"],
    "Lecture 24 -- Last Lecture.pdf": ["现实中的 AI 治理涉及隐私、偏见、可靠性、版权、就业影响和安全边界。", "技术评估应同时关注能力、失败模式、成本和人类监督机制。"],
}


def discover_files() -> list[Path]:
    return sorted(
        p for p in ROOT.iterdir()
        if p.is_file() and p.suffix.lower() in SUPPORTED
        and p.name.lower().startswith("lecture")
    )


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        try:
            import fitz
        except ImportError:
            try:
                from pypdf import PdfReader
            except ImportError as exc:
                raise RuntimeError("缺少 PDF 读取库 PyMuPDF 或 pypdf") from exc
            return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
        with fitz.open(path) as pdf:
            return "\n".join(page.get_text() for page in pdf)
    if suffix == ".pptx":
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("缺少 python-pptx") from exc
        chunks = []
        for slide in Presentation(path).slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    chunks.append(shape.text)
        return "\n".join(chunks)
    if suffix == ".docx":
        try:
            import docx2txt
        except ImportError as exc:
            raise RuntimeError("缺少 docx2txt") from exc
        return docx2txt.process(str(path))
    return path.read_text(encoding="utf-8", errors="replace")


def set_cell_font(run, size: float | None = None, bold: bool | None = None):
    run.font.name = FONT_FAMILY
    run._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_FAMILY)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.bold = bold


def configure_style(style, size: float, bold: bool = False, color=None,
                    before: float = 0, after: float = 0, line: float = 1.0):
    style.font.name = FONT_FAMILY
    style._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_FAMILY)
    style.font.size = Pt(size)
    style.font.bold = bold
    if color:
        style.font.color.rgb = color
    fmt = style.paragraph_format
    fmt.space_before = Pt(before)
    fmt.space_after = Pt(after)
    fmt.line_spacing = line


def set_three_columns(section):
    sect_pr = section._sectPr
    cols = sect_pr.find(qn("w:cols"))
    if cols is None:
        cols = OxmlElement("w:cols")
        sect_pr.append(cols)
    cols.set(qn("w:num"), "3")
    cols.set(qn("w:space"), "360")


def add_page_number(paragraph):
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run("EC508 NLP  |  ")
    set_cell_font(run, BODY_SIZE)
    fld_begin = OxmlElement("w:fldChar")
    fld_begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = " PAGE "
    fld_end = OxmlElement("w:fldChar")
    fld_end.set(qn("w:fldCharType"), "end")
    run._r.append(fld_begin)
    run._r.append(instr)
    run._r.append(fld_end)


def add_bullet(doc: Document, text: str):
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.left_indent = Inches(0.16)
    p.paragraph_format.first_line_indent = Inches(-0.12)
    p.paragraph_format.space_after = Pt(1.5)
    p.paragraph_format.line_spacing = 1.0
    run = p.add_run(text)
    set_cell_font(run, BODY_SIZE)


def add_term(doc: Document, term: str, explanation: str):
    p = doc.add_paragraph(style="Normal")
    p.paragraph_format.space_after = Pt(1)
    run = p.add_run(f"{term}: ")
    set_cell_font(run, BODY_SIZE, True)
    run = p.add_run(explanation)
    set_cell_font(run, BODY_SIZE)


def keyword_fallback(path: Path, text: str) -> dict:
    cleaned = re.sub(r"\s+", " ", text).strip()
    sample = cleaned[:1000] if cleaned else "未提取到可用文本。"
    terms = []
    for term in [
        "language model", "n-gram", "smoothing", "perplexity", "bag-of-words",
        "TF-IDF", "cosine similarity", "PCA", "word embedding", "word2vec",
        "BERT", "Transformer",
    ]:
        if term.lower() in cleaned.lower():
            terms.append((term, "课件中出现的重要 NLP 术语，建议结合原课件复习。"))
    return {
        "title": f"{path.stem}: 自动提取笔记",
        "overview": f"该课件没有预置的结构化摘要。以下内容根据提取文本生成，用于辅助复习：{sample}",
        "points": ["请结合原课件检查自动提取文本中的定义、示例与章节结构。"],
        "terms": terms or [("NLP", "自然语言处理。")],
        "methods": ["未检测到预置公式，请结合原课件复习。"],
        "review": ["复习课件中的标题、定义、示例和方法流程。"],
    }


def add_section(doc: Document, heading: str, items: Iterable[str]):
    doc.add_heading(heading, level=2)
    for item in items:
        add_bullet(doc, item)


def build_document(files: list[Path], extracted: dict[str, str], failed: list[tuple[str, str]]):
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.55)
    section.bottom_margin = Inches(0.55)
    section.left_margin = Inches(0.55)
    section.right_margin = Inches(0.55)
    section.header_distance = Inches(0.25)
    section.footer_distance = Inches(0.25)
    set_three_columns(section)

    styles = doc.styles
    configure_style(styles["Normal"], BODY_SIZE, after=1.5, line=1.0)
    configure_style(styles["Title"], BODY_SIZE, True, ACCENT, after=4)
    configure_style(styles["Heading 1"], BODY_SIZE, True, ACCENT, before=6, after=3)
    configure_style(styles["Heading 2"], BODY_SIZE, True, ACCENT, before=4, after=1.5)
    configure_style(styles["List Bullet"], BODY_SIZE, after=1.5, line=1.0)
    if "Compact Meta" not in [s.name for s in styles]:
        meta = styles.add_style("Compact Meta", WD_STYLE_TYPE.PARAGRAPH)
        configure_style(meta, BODY_SIZE, after=1)

    add_page_number(section.footer.paragraphs[0])

    title = doc.add_paragraph(style="Title")
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("EC508 NLP 课程读书笔记")
    set_cell_font(run, BODY_SIZE, True)
    subtitle = doc.add_paragraph(style="Compact Meta")
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run("中文复习版 | 自动汇总本地 Lecture 课件 | 三栏紧凑排版")
    set_cell_font(run, BODY_SIZE)
    intro = doc.add_paragraph(style="Normal")
    intro.paragraph_format.space_after = Pt(3)
    run = intro.add_run(
        "使用说明：每讲先从需要解决的问题进入，再沿着方法演进梳理核心知识。"
        "术语、公式、考试重点和现实世界应用放在故事线之后，便于先理解动机，再记忆细节。"
        "专业术语保留英文原文，便于对应课件、作业与考试。"
    )
    set_cell_font(run, BODY_SIZE)

    for index, path in enumerate(files):
        if index > 0:
            doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
        note = LECTURE_NOTES.get(path.name) or keyword_fallback(path, extracted.get(path.name, ""))
        doc.add_heading(note["title"], level=1)
        meta = doc.add_paragraph(style="Compact Meta")
        run = meta.add_run(f"文件名：{path.name}")
        set_cell_font(run, BODY_SIZE, True)

        doc.add_heading("引子：问题从哪里来", level=2)
        p = doc.add_paragraph(style="Normal")
        run = p.add_run(NARRATIVE_INTROS.get(path.name, note["overview"]))
        set_cell_font(run, BODY_SIZE)

        doc.add_heading("主题概述", level=2)
        p = doc.add_paragraph(style="Normal")
        run = p.add_run(note["overview"])
        set_cell_font(run, BODY_SIZE)

        add_section(doc, "解决问题的过程：核心知识点", note["points"])

        doc.add_heading("重要术语表", level=2)
        for term, explanation in note["terms"]:
            add_term(doc, term, explanation)

        add_section(doc, "公式或方法总结", note["methods"])
        add_section(doc, "适合复习 / 考试的重点", note["review"])
        add_section(
            doc,
            "额外知识补充：现实世界中的应用",
            REAL_WORLD_APPLICATIONS.get(
                path.name,
                ["该知识点可作为分析真实 NLP 系统的基础，建议结合具体任务评估输入、模型、指标和部署约束。"],
            ),
        )

    doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
    doc.add_heading("未能处理的文件", level=1)
    if failed:
        for name, reason in failed:
            add_bullet(doc, f"{name}: {reason}")
    else:
        p = doc.add_paragraph(style="Normal")
        run = p.add_run("无。所有符合条件的 Lecture / lectures 课程文件均已成功读取并处理。")
        set_cell_font(run, BODY_SIZE)

    props = doc.core_properties
    props.title = "EC508 NLP 课程读书笔记"
    props.subject = "中文复习笔记"
    props.author = "EC508 NLP Notes Generator"
    props.keywords = "NLP, lecture notes, Chinese, EC508"
    try:
        doc.save(OUTPUT)
        return OUTPUT
    except PermissionError:
        fallback = ROOT / "EC508_NLP_lecture_notes_zh_updated.docx"
        doc.save(fallback)
        print(f"目标文件正被占用，已改为生成：{fallback}")
        return fallback


def main() -> int:
    files = discover_files()
    print(f"发现 {len(files)} 个 Lecture / lectures 课程文件：")
    for path in files:
        print(f"  - {path.name}")
    if not files:
        print("未发现可处理文件。", file=sys.stderr)
        return 1

    extracted: dict[str, str] = {}
    failed: list[tuple[str, str]] = []
    usable: list[Path] = []
    for path in files:
        try:
            text = extract_text(path)
            if not text.strip():
                raise RuntimeError("未提取到可用文本")
            extracted[path.name] = text
            usable.append(path)
            print(f"已读取：{path.name}（{len(text)} 字符）")
        except Exception as exc:
            failed.append((path.name, str(exc)))
            print(f"读取失败：{path.name}（{exc}）")

    saved_to = build_document(usable, extracted, failed)
    print(f"已生成：{saved_to}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
